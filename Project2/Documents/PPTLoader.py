
from pptx import Presentation

# Utils/Document.py
from dataclasses import dataclass, field
from typing import Dict, Any, Optional, List

@dataclass
class Document:
    """
    A universal data structure for all document types.
    Works with loaders, chunkers, embeddings, and LLMs.
    """
    page_content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    blob: Optional[bytes] = None       # raw binary (for images)
    resources: Optional[List[str]] = field(default_factory=list)  # extracted file paths, images, etc.
    id: Optional[str] = None
    
    def __repr__(self):
        preview = self.page_content[:80].replace("\n", " ") + ("..." if len(self.page_content) > 80 else "")
        return f"<Document id={self.id} source={self.metadata.get('source','')} content='{preview}'>"

class LoadPPTX:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list:
        prs = Presentation(self.file_path)
        documents = []

        for i, slide in enumerate(prs.slides):
            text = ""
            images = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    images.append(image.blob)

            documents.append(Document(
                page_content=text.strip(),
                metadata={"source": self.file_path, "slide": i + 1},
                resources=images  # ✅ fixed field name
            ))
        self._documents = documents
        return documents
    def __iter__(self):
        if self._documents is None:
            self.load()
        return iter(self._documents)
